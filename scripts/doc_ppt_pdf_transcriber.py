import sys
import os
import json
import tempfile

def find_python():
    candidates = [
        r"C:\Users\corte\Documents\projects NOT DELETE\lex memoria\venv\Scripts\python.exe",
        sys.executable
    ]
    for c in candidates:
        if os.path.exists(c):
            return c
    return sys.executable

def extract_from_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    slides_text = []
    for idx, slide in enumerate(prs.slides):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                lines.append(f"Notes: {note}")
        
        slide_content = "\n".join(lines) if lines else "[Slide contains diagrams or images only]"
        slides_text.append(f"=== Slide {idx + 1} ===\n{slide_content}")
    return "\n\n".join(slides_text)

def extract_from_pdf(file_path):
    try:
        import fitz
    except ImportError:
        import pymupdf as fitz
    
    doc = fitz.open(file_path)
    pages_text = []
    for idx, page in enumerate(doc):
        text = page.get_text("text").strip()
        # If very low text, check for images/diagrams
        if len(text) < 50:
            image_list = page.get_images(full=True)
            if image_list:
                text += f"\n[Page contains {len(image_list)} figures/diagrams]"
        if not text:
            text = f"[Page {idx + 1}]"
        pages_text.append(f"=== Page {idx + 1} ===\n{text}")
    doc.close()
    return "\n\n".join(pages_text)

def main():
    if len(sys.argv) < 2:
        print("Usage: doc_ppt_pdf_transcriber.py <file_path>", file=sys.stderr)
        sys.exit(1)
        
    file_path = sys.argv[1]
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext in [".pptx", ".ppt"]:
        extracted = extract_from_pptx(file_path)
    elif ext == ".pdf":
        extracted = extract_from_pdf(file_path)
    else:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            extracted = f.read()
            
    print(extracted)

if __name__ == "__main__":
    main()
